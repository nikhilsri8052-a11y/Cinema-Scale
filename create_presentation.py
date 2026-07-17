"""
Cinema-Scale Group Presentation Generator
==========================================
Generates a professional .pptx presentation for the Cinema-Scale project,
divided into 5 sections so each group member can present their assigned part.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
DARK_BG       = RGBColor(0x0F, 0x0F, 0x1A)   # Deep navy-black
ACCENT_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # Vibrant purple
ACCENT_CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # Bright cyan
ACCENT_PINK   = RGBColor(0xEC, 0x48, 0x99)   # Hot pink
ACCENT_GREEN  = RGBColor(0x10, 0xB9, 0x81)   # Emerald green
ACCENT_AMBER  = RGBColor(0xF5, 0x9E, 0x0B)   # Warm amber
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY      = RGBColor(0x94, 0xA3, 0xB8)
DARK_CARD     = RGBColor(0x1E, 0x1E, 0x2E)   # Card background

# Member accent colors (one per presenter)
MEMBER_COLORS = [ACCENT_PURPLE, ACCENT_CYAN, ACCENT_PINK, ACCENT_GREEN, ACCENT_AMBER]

# ─── Presentation Setup ─────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ─── Helper Functions ────────────────────────────────────────────────────────

def add_solid_bg(slide, color=DARK_BG):
    """Set solid dark background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.rotation = 0
    sf = shape.fill
    if fill_color:
        sf.solid()
        sf.fore_color.rgb = fill_color
    else:
        sf.background()
    ln = shape.line
    if border_color:
        ln.color.rgb = border_color
        ln.width = border_width
    else:
        ln.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_frame(slide, left, top, width, height, items, font_size=16,
                     color=LIGHT_GRAY, accent=ACCENT_PURPLE, font_name="Calibri"):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.space_after = Pt(4)

        # Bullet symbol
        run_bullet = p.add_run()
        run_bullet.text = "▸ "
        run_bullet.font.size = Pt(font_size)
        run_bullet.font.color.rgb = accent
        run_bullet.font.name = font_name
        run_bullet.font.bold = True

        # Item text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name
    return txBox


def add_section_divider(slide, member_name, section_title, member_color, section_num):
    """Add a section divider slide for each presenter."""
    add_solid_bg(slide)

    # Large section number
    add_text_box(slide, Inches(1), Inches(0.8), Inches(3), Inches(2),
                 f"0{section_num}", font_size=96, color=member_color,
                 bold=True, font_name="Calibri Light")

    # Decorative line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(3), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = member_color
    line.line.fill.background()

    # Section title
    add_text_box(slide, Inches(1), Inches(3.5), Inches(10), Inches(1.2),
                 section_title, font_size=40, color=WHITE, bold=True, font_name="Calibri")

    # Presenter name
    add_text_box(slide, Inches(1), Inches(4.8), Inches(10), Inches(0.8),
                 f"Presented by  {member_name}", font_size=22, color=MID_GRAY,
                 font_name="Calibri")

    # Accent circle decoration (top-right)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2), Inches(2))
    circle.fill.solid()
    circle.fill.fore_color.rgb = member_color
    # Make it semi-transparent by setting alpha
    circle.fill.fore_color.brightness = 0.6
    circle.line.fill.background()


def add_content_slide(slide, title, bullets, accent_color, sub_title=None):
    """Standard content slide with title and bullet points."""
    add_solid_bg(slide)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                 title, font_size=34, color=WHITE, bold=True, font_name="Calibri")

    # Subtle underline
    uline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(2.5), Pt(3))
    uline.fill.solid()
    uline.fill.fore_color.rgb = accent_color
    uline.line.fill.background()

    y_start = 1.6
    if sub_title:
        add_text_box(slide, Inches(0.8), Inches(y_start), Inches(11), Inches(0.6),
                     sub_title, font_size=18, color=MID_GRAY, font_name="Calibri")
        y_start += 0.6

    # Bullet content
    add_bullet_frame(slide, Inches(1.0), Inches(y_start), Inches(10.5), Inches(5),
                     bullets, font_size=18, color=LIGHT_GRAY, accent=accent_color)


def add_two_column_slide(slide, title, left_title, left_items, right_title, right_items, accent_color):
    """Slide with two columns of content."""
    add_solid_bg(slide)

    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Pt(6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()

    # Title
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                 title, font_size=34, color=WHITE, bold=True, font_name="Calibri")

    uline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(2.5), Pt(3))
    uline.fill.solid()
    uline.fill.fore_color.rgb = accent_color
    uline.line.fill.background()

    # Left column card
    add_shape(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), fill_color=DARK_CARD, border_color=accent_color, border_width=Pt(1))
    add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.5),
                 left_title, font_size=22, color=accent_color, bold=True)
    add_bullet_frame(slide, Inches(1.1), Inches(2.6), Inches(5), Inches(4),
                     left_items, font_size=16, accent=accent_color)

    # Right column card
    add_shape(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), fill_color=DARK_CARD, border_color=accent_color, border_width=Pt(1))
    add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.5),
                 right_title, font_size=22, color=accent_color, bold=True)
    add_bullet_frame(slide, Inches(7.1), Inches(2.6), Inches(5), Inches(4),
                     right_items, font_size=16, accent=accent_color)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_solid_bg(slide)

# Large decorative circle (background)
deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-1), Inches(6), Inches(6))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_PURPLE
deco.fill.fore_color.brightness = 0.7
deco.line.fill.background()

deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.5), Inches(5), Inches(5))
deco2.fill.solid()
deco2.fill.fore_color.rgb = ACCENT_CYAN
deco2.fill.fore_color.brightness = 0.7
deco2.line.fill.background()

# Emoji + Title
add_text_box(slide, Inches(1), Inches(1.2), Inches(10), Inches(0.8),
             "🎬", font_size=54, alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(2.2), Inches(10), Inches(1.2),
             "Cinema-Scale", font_size=60, color=WHITE, bold=True,
             font_name="Calibri", alignment=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1), Inches(3.4), Inches(10), Inches(0.7),
             "An Intelligent Movie Browsing & Content-Based Recommendation Platform",
             font_size=22, color=LIGHT_GRAY, font_name="Calibri", alignment=PP_ALIGN.LEFT)

# Separator
sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.3), Inches(4), Pt(3))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_PURPLE
sep.line.fill.background()

# Team members
members = [
    "Nikhil Srivastava", "Kritika Saini", "Deepika Kattal",
    "Diya Sharma", "Satkar Singh"
]
for i, name in enumerate(members):
    add_text_box(slide, Inches(1), Inches(4.6 + i * 0.4), Inches(6), Inches(0.4),
                 f"●  {name}", font_size=16, color=MEMBER_COLORS[i], font_name="Calibri")

add_text_box(slide, Inches(1), Inches(6.7), Inches(10), Inches(0.4),
             "IIT Madras  •  Group Project Presentation", font_size=14,
             color=MID_GRAY, font_name="Calibri")


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — AGENDA / TABLE OF CONTENTS
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
             "Presentation Agenda", font_size=38, color=WHITE, bold=True)

uline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35), Inches(3), Pt(3))
uline.fill.solid()
uline.fill.fore_color.rgb = ACCENT_PURPLE
uline.line.fill.background()

agenda_items = [
    ("01", "Project Overview & Architecture", "Nikhil Srivastava", ACCENT_PURPLE),
    ("02", "Recommendation Engine & ML Pipeline", "Kritika Saini", ACCENT_CYAN),
    ("03", "Database Design & Data Pipeline", "Deepika Kattal", ACCENT_PINK),
    ("04", "Authentication & Admin System", "Diya Sharma", ACCENT_GREEN),
    ("05", "Frontend Design & User Experience", "Satkar Singh", ACCENT_AMBER),
]

for i, (num, topic, name, color) in enumerate(agenda_items):
    y = 1.8 + i * 1.05
    # Card background
    add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.85),
              fill_color=DARK_CARD, border_color=color, border_width=Pt(1.5))
    # Number
    add_text_box(slide, Inches(1.1), Inches(y + 0.1), Inches(0.8), Inches(0.6),
                 num, font_size=28, color=color, bold=True, font_name="Calibri Light")
    # Topic
    add_text_box(slide, Inches(2.0), Inches(y + 0.05), Inches(6), Inches(0.4),
                 topic, font_size=20, color=WHITE, bold=True, font_name="Calibri")
    # Name
    add_text_box(slide, Inches(2.0), Inches(y + 0.45), Inches(6), Inches(0.35),
                 name, font_size=14, color=MID_GRAY, font_name="Calibri")
    # Colored accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.3), Inches(y + 0.25), Inches(0.35), Inches(0.35))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 1 — NIKHIL SRIVASTAVA: Project Overview & Architecture
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "Nikhil Srivastava", "Project Overview & Architecture", ACCENT_PURPLE, 1)

# Slide 1.1 — What is Cinema-Scale?
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "What is Cinema-Scale?", [
    "An interactive, full-stack movie browsing & recommendation platform",
    "Built with Python/Flask backend and Machine Learning algorithms",
    "Content-based filtering using TF-IDF Vectorization & Cosine Similarity",
    "Dynamically enriches movie metadata from TMDB API & Wikipedia CDN",
    "Supports user authentication, reviews, watchlists, and admin controls",
    "Processes 45,000+ movies from a comprehensive CSV dataset",
], ACCENT_PURPLE, sub_title="A full-stack movie platform powered by machine learning")

# Slide 1.2 — Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_column_slide(slide, "Technology Stack",
    "Backend & ML", [
        "Flask — Lightweight Python web framework",
        "SQLAlchemy ORM — Database abstraction layer",
        "SQLite / PostgreSQL — Relational database engines",
        "Scikit-learn — TF-IDF & Cosine Similarity",
        "Pandas & NumPy — Data processing",
        "difflib — Fuzzy string matching",
    ],
    "Frontend & APIs", [
        "HTML5 + Vanilla CSS3 — Custom glassmorphism UI",
        "JavaScript — AJAX pagination & dynamic loading",
        "TMDB API — Live movie metadata & poster CDN",
        "Wikipedia Media API — Poster fallback source",
        "Responsive design — Mobile-first approach",
        "Interactive carousels & hover animations",
    ], ACCENT_PURPLE)

# Slide 1.3 — System Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "System Architecture", [
    "Flask Web Client sends HTTP requests to Controller Blueprints",
    "Controller queries the SQLAlchemy Database Cache for movie data",
    "Controller checks the In-Memory Similarity Engine for recommendations",
    "Similarity Engine performs index lookups against the precomputed TF-IDF Matrix",
    "If metadata is missing, the TMDB/Wikipedia Scraper Helper is triggered",
    "Scraped & enriched metadata is cached back into the database",
    "Eager initialization: ML model trains at server startup for O(1) lookups",
], ACCENT_PURPLE, sub_title="Request flow: Client → Controller → ML Engine → Database → API Enrichment")

# Slide 1.4 — Initialization & Startup Flow
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Application Initialization Flow", [
    "Step 1: Database setup — dynamically selects SQLite (local) or PostgreSQL (production)",
    "Step 2: Model training — calls load_and_train_model('movies.csv') at boot",
    "Step 3: Builds TF-IDF feature matrix & precomputes Cosine Similarity matrix in-memory",
    "Step 4: Database seeding — runs seed_database_from_csv() to populate tables in chunks",
    "Step 5: Context processors inject session data (current_user, is_logged_in, is_admin)",
    "Step 6: Registers all Flask Blueprints (auth, main, interactions, admin, api)",
], ACCENT_PURPLE, sub_title="What happens when the server starts (app.py)")


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 2 — KRITIKA SAINI: Recommendation Engine & ML Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "Kritika Saini", "Recommendation Engine & ML Pipeline", ACCENT_CYAN, 2)

# Slide 2.1 — Content-Based Filtering Approach
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Content-Based Filtering Approach", [
    "Recommends movies based on metadata similarity, not user behavior",
    "Analyzes genres, keywords, tagline, cast, and director features",
    "Combines all text features into a single metadata 'document' per movie",
    "Uses NLP (Natural Language Processing) techniques for vectorization",
    "Fully deterministic: same input always yields the same recommendations",
    "No cold-start problem — works even for new users with zero history",
], ACCENT_CYAN, sub_title="How we match similar movies using metadata")

# Slide 2.2 — Feature Engineering
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Feature Engineering", [
    "5 key features combined: genres + keywords + tagline + cast + director",
    "Example: 'Action Adventure space travel Matthew McConaughey Christopher Nolan'",
    "Missing values handled with empty string fallbacks to prevent NaN errors",
    "Combined features create a rich text representation for each movie",
    "This 'bag of words' approach captures the essence of each film's identity",
    "Features are normalized and cleaned before vectorization step",
], ACCENT_CYAN, sub_title="Combining movie attributes into a unified text representation")

# Slide 2.3 — TF-IDF Vectorization
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "TF-IDF Vectorization", [
    "TF-IDF = Term Frequency × Inverse Document Frequency",
    "Term Frequency (TF): How often a keyword appears in one movie's metadata",
    "Inverse Document Frequency (IDF): Penalizes common words like 'the', 'movie'",
    "Boosts rare, descriptive terms: 'cyberpunk', 'inception', 'time-travel'",
    "Scikit-learn's TfidfVectorizer transforms text into high-dimensional vectors",
    "Each movie becomes a numerical vector in word-space for comparison",
], ACCENT_CYAN, sub_title="Converting text metadata into numerical feature vectors")

# Slide 2.4 — Cosine Similarity & Retrieval
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_column_slide(slide, "Cosine Similarity & Retrieval",
    "Cosine Similarity", [
        "Measures angle between two TF-IDF vectors",
        "Formula: cos(A,B) = (A·B) / (||A|| × ||B||)",
        "Score range: 0 (no similarity) to 1 (identical)",
        "Generates N×N similarity matrix at startup",
        "O(1) lookup time after precomputation",
    ],
    "Fuzzy Matching & Retrieval", [
        "difflib.get_close_matches handles user typos",
        "Example: 'Interstelar' → 'Interstellar'",
        "Finds closest title in the dataset index",
        "Retrieves similarity row, sorts descending",
        "Returns top N recommendations (default: 10)",
    ], ACCENT_CYAN)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 3 — DEEPIKA KATTAL: Database Design & Data Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "Deepika Kattal", "Database Design & Data Pipeline", ACCENT_PINK, 3)

# Slide 3.1 — Database Architecture
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Database Architecture Overview", [
    "Built on SQLAlchemy ORM mapped to SQLite (local) or PostgreSQL (production)",
    "5 core tables: users, movies, likes, reviews, saved_movies",
    "Movies table stores 18+ columns of deep metadata including TMDB fields",
    "Join tables (likes, saved_movies) enforce UniqueConstraints to prevent duplicates",
    "Relationships defined with foreign keys and back-references for efficient querying",
    "Timestamps on all user actions for analytics and audit trails",
], ACCENT_PINK, sub_title="SQLAlchemy ORM with SQLite/PostgreSQL backend")

# Slide 3.2 — Entity Relationship Diagram
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Entity Relationship Diagram (ERD)", [
    "Users (1) ──── has many ────→ Likes (M) ←──── belongs to ──── Movies (1)",
    "Users (1) ──── has many ────→ Reviews (M) ←── belongs to ──── Movies (1)",
    "Users (1) ──── has many ────→ Saved Movies (M) ←─ belongs to ─ Movies (1)",
    "Users table: id, username, email, password_hash, is_admin, created_at",
    "Movies table: id, tmdb_id, title, genre, year, plot, rating, poster_path, runtime, budget, revenue, etc.",
    "Reviews: rating (1-10 scale) + comment text + timestamp",
], ACCENT_PINK, sub_title="5 interconnected tables with foreign key relationships")

# Slide 3.3 — Data Seeding Pipeline
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Data Seeding & Enrichment Pipeline", [
    "Source dataset: movies.csv with 45,000+ movie records",
    "Chunk-based seeding: loads data in batches to prevent memory overflow",
    "On-demand metadata enrichment via TMDB API (budget, runtime, revenue)",
    "Wikipedia Media API fallback when TMDB is rate-limited or unavailable",
    "Poster paths cached in database to avoid repeated API calls",
    "populate_posters.py script for batch pre-populating poster URLs",
], ACCENT_PINK, sub_title="From CSV to enriched database with live API integration")

# Slide 3.4 — Key Design Decisions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_column_slide(slide, "Database Design Decisions",
    "Performance Optimizations", [
        "Eager model training at startup (not per-request)",
        "Database caching prevents repeated API calls",
        "Chunk-based seeding avoids memory overflow",
        "Indexed foreign keys for fast JOIN queries",
        "Lazy poster loading — fetched only when viewed",
    ],
    "Data Integrity", [
        "UniqueConstraint on likes (user_id, movie_id)",
        "UniqueConstraint on saved_movies (user_id, movie_id)",
        "Unique email constraint on users table",
        "NOT NULL constraints on critical fields",
        "Password hashing for secure authentication",
    ], ACCENT_PINK)


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 4 — DIYA SHARMA: Authentication & Admin System
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "Diya Sharma", "Authentication & Admin System", ACCENT_GREEN, 4)

# Slide 4.1 — User Authentication
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "User Authentication System", [
    "Secure registration with email uniqueness validation",
    "Password hashing using Werkzeug's generate_password_hash",
    "Session-based login/logout management via Flask sessions",
    "Context processors inject user state into every page template",
    "Role-based access: regular users vs. administrators",
    "Protected routes: watchlist, reviews, and profile require login",
], ACCENT_GREEN, sub_title="Secure registration, login, and session management")

# Slide 4.2 — User Interactions
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_column_slide(slide, "User Interactions & Features",
    "Engagement Features", [
        "Like/unlike toggle on any movie page",
        "Save to watchlist for later viewing",
        "Write text reviews with 1-10 rating scale",
        "View personal profile with liked & saved movies",
        "Dynamic average rating calculation",
    ],
    "Blueprint Architecture", [
        "auth.py — Registration, login, logout routes",
        "interactions.py — Likes, reviews, saves handling",
        "main.py — Landing, details, category, profile pages",
        "api.py — AJAX paginated movie loading",
        "admin.py — Administrative controls & dashboard",
    ], ACCENT_GREEN)

# Slide 4.3 — Admin Dashboard
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Administrative Dashboard", [
    "Accessible only to users with is_admin = True flag",
    "System telemetry: total movies, users, and reviews counts",
    "Add new movies manually via web form interface",
    "Delete movies, user accounts, or violating reviews",
    "Monitor and manage community content moderation",
    "Protected by role-checking middleware on /admin route",
], ACCENT_GREEN, sub_title="Full CRUD operations and system monitoring for admins")

# Slide 4.4 — Security Measures
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Security & Access Control", [
    "Passwords never stored in plaintext — always hashed",
    "Session-based authentication with server-side state",
    "CSRF protection through Flask's built-in mechanisms",
    "Input validation on all user-facing forms",
    "Admin-only route guards prevent unauthorized access",
    "Graceful error handling for invalid login attempts",
], ACCENT_GREEN, sub_title="Multi-layered security across the application")


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION 5 — SATKAR SINGH: Frontend Design & User Experience
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, "Satkar Singh", "Frontend Design & User Experience", ACCENT_AMBER, 5)

# Slide 5.1 — UI/UX Design Philosophy
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "UI/UX Design Philosophy", [
    "Modern glassmorphism aesthetic with frosted-glass card effects",
    "Dark theme with vibrant accent colors for premium feel",
    "Responsive grid layouts that adapt to all screen sizes",
    "Smooth CSS transitions and hover micro-animations throughout",
    "Interactive carousels for movie browsing and recommendations",
    "Consistent visual language across all pages and components",
], ACCENT_AMBER, sub_title="Creating a premium, immersive movie browsing experience")

# Slide 5.2 — Key Frontend Pages
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_column_slide(slide, "Key Frontend Pages & Components",
    "Core Pages", [
        "Landing page with hero section & trending carousel",
        "Movie details page with full metadata display",
        "Category browsing with genre-based filtering",
        "User profile with saved movies & reviews",
        "Admin dashboard with telemetry cards",
    ],
    "Interactive Components", [
        "AJAX infinite scroll pagination (no page reloads)",
        "Dynamic like/save toggle buttons",
        "Star rating input widget (1-10 scale)",
        "Glassmorphic movie cards with hover effects",
        "Responsive navigation with session-aware links",
    ], ACCENT_AMBER)

# Slide 5.3 — CSS & Styling Details
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Styling & Visual Design", [
    "Pure Vanilla CSS3 — no CSS frameworks, fully custom-built",
    "Glassmorphism: backdrop-filter blur + transparency on cards",
    "Custom color palette with dark backgrounds and neon accents",
    "Smooth transitions on all interactive elements (0.3s ease)",
    "Responsive media queries for mobile, tablet, and desktop",
    "Google Fonts integration for modern, clean typography",
], ACCENT_AMBER, sub_title="Hand-crafted CSS with modern design techniques")

# Slide 5.4 — JavaScript & Dynamic Features
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "JavaScript & Dynamic Features", [
    "AJAX-powered pagination — loads movies without full page reload",
    "Fetch API calls to /api/movies endpoint for lazy loading",
    "Dynamic DOM manipulation for real-time like/save state updates",
    "Smooth scroll animations for improved navigation feel",
    "Error handling with user-friendly toast notifications",
    "Optimized rendering: only loads visible movie cards on demand",
], ACCENT_AMBER, sub_title="Client-side interactivity and performance optimizations")


# ═══════════════════════════════════════════════════════════════════════════════
#  CLOSING SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

# Future Scope slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide, "Future Scope & Enhancements", [
    "Collaborative Filtering — incorporate user behavior patterns alongside content similarity",
    "Deep Learning Models — use neural embeddings (Word2Vec / BERT) for richer representations",
    "Real-time Streaming — integrate WebSockets for live updates and notifications",
    "Mobile App — build native iOS/Android apps with React Native",
    "Cloud Deployment — containerize with Docker and deploy on AWS/GCP",
    "A/B Testing — experiment with recommendation algorithms at scale",
], ACCENT_PURPLE, sub_title="Where Cinema-Scale goes from here")

# Thank You slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_solid_bg(slide)

# Decorative elements
deco = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(-1.5), Inches(6), Inches(6))
deco.fill.solid()
deco.fill.fore_color.rgb = ACCENT_PURPLE
deco.fill.fore_color.brightness = 0.7
deco.line.fill.background()

deco2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(5), Inches(5))
deco2.fill.solid()
deco2.fill.fore_color.rgb = ACCENT_CYAN
deco2.fill.fore_color.brightness = 0.7
deco2.line.fill.background()

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "Thank You!", font_size=64, color=WHITE, bold=True,
             font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
             "Cinema-Scale  •  Movie Recommendation Platform", font_size=24,
             color=LIGHT_GRAY, font_name="Calibri", alignment=PP_ALIGN.CENTER)

sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.2), Inches(3.333), Pt(3))
sep.fill.solid()
sep.fill.fore_color.rgb = ACCENT_PURPLE
sep.line.fill.background()

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Questions & Discussion", font_size=22, color=ACCENT_PURPLE,
             bold=True, font_name="Calibri", alignment=PP_ALIGN.CENTER)

# Team credits
for i, name in enumerate(members):
    col = i % 5
    x = 1.0 + col * 2.3
    add_text_box(slide, Inches(x), Inches(5.5), Inches(2.2), Inches(0.4),
                 f"●  {name}", font_size=13, color=MEMBER_COLORS[i],
                 font_name="Calibri", alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
             "github.com/nikhilsri8052-a11y/Cinema-Scale", font_size=14,
             color=MID_GRAY, font_name="Calibri", alignment=PP_ALIGN.CENTER)


# ─── Save Presentation ──────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "Cinema-Scale_Presentation.pptx")
prs.save(output_path)
print(f"[OK] Presentation saved successfully to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print(f"\nSlide breakdown:")
print(f"  - Slide 1:      Title Slide (all members)")
print(f"  - Slide 2:      Agenda / Table of Contents")
print(f"  - Slides 3-7:   Section 1 -- Nikhil Srivastava (Project Overview & Architecture)")
print(f"  - Slides 8-12:  Section 2 -- Kritika Saini (Recommendation Engine & ML Pipeline)")
print(f"  - Slides 13-17: Section 3 -- Deepika Kattal (Database Design & Data Pipeline)")
print(f"  - Slides 18-22: Section 4 -- Diya Sharma (Authentication & Admin System)")
print(f"  - Slides 23-27: Section 5 -- Satkar Singh (Frontend Design & User Experience)")
print(f"  - Slide 28:     Future Scope")
print(f"  - Slide 29:     Thank You & Q&A")
